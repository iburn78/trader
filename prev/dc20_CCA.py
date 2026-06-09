#%% 
# CCA: Company Classification Analysis
from trader.tools.ca_tools import get_score_trend
import pickle
from openai import OpenAI
import json
import re
from datetime import datetime, timedelta
import FinanceDataReader as fdr
import io

# ----------------------------------------------------------
# Post processing and PPT Generation
# ----------------------------------------------------------

###_ ------------------------------------------------
###_ below outshare_db should be fixed - inefficient and no longer maintains outshare_db
###_ ------------------------------------------------

from trader.tools.dc_tools import get_main_financial_reports_db, get_quarterly_data, get_price_db, get_outshare_db, prev_quarter_str
from trader.tools.dc_tools import plot_company_financial_summary
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches

fr_db = get_main_financial_reports_db()
pr_db = get_price_db()
outshare_DB = get_outshare_db()
target_account = 'net_income'

# this is a different version of L4_addition function from anaysis project
def L4_rolling_addition(fh, target_account=target_account):
    # Add L4 key account to the financial report database.
    new_row = {'account':'L4_'+target_account }
    target_row = fh.loc[target_account]

    for i in range(3, len(fh.columns)):
        previous_4_quarters = fh.columns[i-3:i+1]
        rolling_sum = target_row[previous_4_quarters].sum()
        new_row[fh.columns[i]] = rolling_sum

    new_row_df = pd.DataFrame([new_row]).set_index('account')
    fh = pd.concat([fh, new_row_df])
    return fh 

# get price, marcap, PER, PBR
def get_market_performance_db(code, fr_db=fr_db, pr_db=pr_db, outshare_DB=outshare_DB, target_account=target_account, MAX_QUARTERS=40): # e.g., 10 years or use MAX_QUARTERS in qa_db generation
    fh = get_quarterly_data(code, fr_db, native=True)
    fh = L4_rolling_addition(fh, target_account)
    try:
        init_loc = fh.columns.get_loc(fh.loc['L4_'+ target_account].first_valid_index())
        start_loc = max(max(len(fh.columns) - MAX_QUARTERS, 0), init_loc)
    except:
        start_loc = 0
    start_day  = pd.Period(fh.columns[start_loc][:5].replace('_', ''), freq='Q').start_time.strftime('%Y-%m-%d')

    # market performance
    mp_db = pd.DataFrame()
    mp_db['price'] = pr_db[code].loc[pr_db.index >= start_day]
    mp_db.index = pd.to_datetime(mp_db.index)
    shares = outshare_DB[code].iloc[-1]
    mp_db['marcap'] = mp_db['price'] * shares 
    r_ = fh.loc['L4_' + target_account].dropna()
    latest_l4 = r_.iloc[-1] if not r_.empty else np.nan
    mp_db['L4'] = mp_db.index.map(lambda d: fh.at['L4_'+target_account, prev_quarter_str(d)] if prev_quarter_str(d) in fh.columns else latest_l4)
    mp_db['PER'] = mp_db['marcap'] / mp_db['L4']

    b_ = fh.loc['equity'].dropna()
    latest_b = b_.iloc[-1] if not b_.empty else np.nan
    mp_db['equity'] = mp_db.index.map(lambda d: fh.at['equity', prev_quarter_str(d)] if prev_quarter_str(d) in fh.columns else latest_b)
    mp_db['PBR'] = mp_db['marcap'] / mp_db['equity']
    mp_db = mp_db.apply(pd.to_numeric, errors='coerce')

    return mp_db


def mp_plot(mp_db, columns=['price', 'PER', 'PBR']):
    fig, axes = plt.subplots(len(columns), 1, figsize=(12, 4 * len(columns)), sharex=True)
    if len(columns) == 1:
        axes = [axes]
    for ax, col in zip(axes, columns):
        mp_db[col].plot(ax=ax, title=col, grid=True, fontsize=12)
    plt.tight_layout()
    img_stream = io.BytesIO()
    plt.savefig(img_stream, format='png', bbox_inches='tight')
    plt.close(fig)
    img_stream.seek(0)
    return img_stream
    # plt.show()
    # plt.savefig("----.png")

def get_late_prices(codelist, days=7):
    start_date = (datetime.today() - timedelta(days=days)).strftime('%Y-%m-%d')
    res = pd.DataFrame()

    for code in codelist:
        df = fdr.DataReader(code, start_date)['Close']
        df.name = code  # name the Series as the stock code
        res = pd.concat([res, df], axis=1)

    res = res.T
    res.columns = pd.to_datetime(res.columns).strftime("%m-%d")
    return res  # transpose to have stock codes as index and dates as columns

def get_codelist_summary(codelist, days=7, df_krx=df_krx):
    last_prices = get_late_prices(codelist, days)
    pr_changes = (last_prices.pct_change(axis=1) * 100).round(2).dropna(axis=1)
    pr_changes.insert(0, 'name', '')
    for code in pr_changes.index:
        pr_changes.loc[code, 'name'] = df_krx.loc[code, "Name"]
        pr_changes.loc[code, 'price'] = last_prices.loc[code, last_prices.columns[-1]]
        pr_changes.loc[code, 'marcap'] = df_krx.loc[code, "Marcap"]/10**8
    pr_changes['marcap'] = pr_changes['marcap'].round().astype(int)
    return pr_changes

def post_process(score_trend, qa_db=qa_db, fr_db=fr_db, pr_db=pr_db, outshare_DB=outshare_DB, df_krx=df_krx, top_N = top_N):
    data_dict = {}
    mp_db_dict = {}
    periods = get_periods(qa_db)

    cls = get_codelist_summary(score_trend.index)
    for code in score_trend.index: 
        print("processing", code)
        mp_db = get_market_performance_db(code, fr_db, pr_db, outshare_DB)
        mp_db_dict[code] = mp_db
        cls.loc[code, 'PER'] = round(mp_db['PER'].iloc[-1], 2)
        rank = str(score_trend.index.get_loc(code) + 1)
        cls.loc[code, 'rank'] = rank 
        txt = score_trend.loc[[code]][periods+['avg']].to_string(index=False)
        cls.loc[code, 'note'] = 'rank:' + rank + ', select:' + score_trend.loc[code, 'selected'] + f', top{str(top_N)}:' + score_trend.loc[code, f'top{str(top_N)}'] + ', MCap:' + str(df_krx.at[code, 'Marcap'] // 10**8) + ', PER:' + str(round(mp_db['PER'].iloc[-1], 2)) + '\n' + txt

    data_dict['codelist_summary'] = cls
    data_dict['mp_db_dict'] = mp_db_dict

    # codelist selection logic ------------------
    PER_limit = 1000
    scodelist = cls.loc[(cls['PER'] > 0) & (cls['PER'] < PER_limit)].index
    data_dict['select_codelist'] = scodelist
    data_dict['select_codelist_summary'] = cls.loc[cls.index.isin(scodelist)].drop('note', axis=1)


    return data_dict

def generate_PPT(data_dict, fr_db=fr_db, pr_db=pr_db, summary_only = False, top_N: int = top_N):

    cd_ = os.path.dirname(os.path.abspath(__file__)) # .
    CCA_folder = 'CCA'
    today_str = pd.Timestamp.today().strftime('%Y-%m-%d_%H%M')
    CCA_template = os.path.join(cd_, CCA_folder, 'CCA_template.pptx')
    CCA_result = os.path.join(cd_, CCA_folder, f'CCA_result_{today_str}.pptx')
    prs = Presentation(CCA_template)

    # Codelist summary page 
    slide = prs.slides.add_slide(prs.slide_layouts[2])
    img_path = styled_df_to_image(data_dict['select_codelist_summary'])
    slide.shapes.add_picture(img_path, 0, Inches(0.1), width=prs.slide_width*0.55)
    os.remove(img_path)

    # Per code pages 
    if not summary_only: 
        if top_N is None: top_N = len(data_dict['select_codelist'])
        for code in data_dict['select_codelist'][:top_N]:
            print("slide generating for", code)
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            img_stream = mp_plot(data_dict['mp_db_dict'][code])
            slide.shapes.add_picture(img_stream, Inches(0.1), Inches(0.5), height=Inches(6))
            for ph in slide.placeholders: 
                if ph.name == 'Text Placeholder 1':
                    ph.text = data_dict['codelist_summary'].loc[code, 'name']
                if ph.name == 'Text Placeholder 2':
                    ph.text = today_str
                if ph.name == 'Text Placeholder 3':
                    ph.text = data_dict['codelist_summary'].loc[code, 'note']

            slide = prs.slides.add_slide(prs.slide_layouts[1])
            ph = next(iter(slide.placeholders)) 
            ph.text = ""
            # ph.text = openai_command(code)

            slide = prs.slides.add_slide(prs.slide_layouts[2])
            img_stream = plot_company_financial_summary(fr_db, pr_db, code, None) 
            slide.shapes.add_picture(img_stream, Inches(0.1), 0, height=prs.slide_height)

            slide = prs.slides.add_slide(prs.slide_layouts[2])
            img_path = gen_data_in_html(code)
            slide.shapes.add_picture(img_path, 0, Inches(0.1), width=prs.slide_width)
            os.remove(img_path)

    prs.save(CCA_result)
    print("PPT generation completed...")

# open ai original
# def openai_command(company_name, conf_file=CONF_FILE):  
#     with open(conf_file, 'r') as json_file:
#         config = json.load(json_file)
#         api_key = config['openai']

#     client = OpenAI(
#         api_key=api_key
#     )
#     content_command = f"{company_name} is a listed company in Korea. Give me brief description (1) what is this company's business, (2) why this company's performance is good for the last quarters, (3) what are the key 3 issues ths company is facing?"
#     format_command = "Answer in Korean language, and use only plain text. Total length of answers should not exceed 500 words."

#     chat_completion = client.chat.completions.create(
#         model="gpt-4o-mini", 
#         messages=[
#             {
#                 "role": "user",
#                 "content": (
#                     content_command + " " + format_command
#                 )
#             }
#         ]
#     )
#     return chat_completion.choices[0].message.content

# perplexity
def openai_command(code, conf_file=CONF_FILE):
    company_name = df_krx.loc[code, "Name"]

    content_command = (
        f"{company_name}({code}) is a listed company in Korea. Provide a detailed and specific explanation of "
        "(1) what this company's core business is, "
        "(2) why the company has shown strong performance in recent quarters, and "
        "(3) what the top 3 key issues the company is currently facing. "
    )

    format_command = (
        "Answer in Korean using plain text. "
        "Avoid vague or high-level generalizations. "
        "The total length must not exceed 700 words. "
    )

    info, prev, date_prev = get_company_issues(code)

    # if the prev gpt response is saved within 7 days, pass
    if datetime.today() - datetime.strptime(date_prev, '%Y-%m-%d') < timedelta(days=7):
        if prev:
            return prev

    if info:
        info_command = (
            "The following is additional information. "
            "Cross-check its validity and use it only if it is important and accurate.\n"
        )
        info = info_command + info

    if prev:
        prev_command = (
            f"The following is your previous response to the identical query about the company on {date_prev}. "
            "Refer to it, but update your answer with any new developments since then.\n"
        )
        prev = prev_command + prev

    extras = "\n\n".join(filter(None, [info, prev]))
    prompt = f"{content_command}\n{format_command}\n\n{extras}"

    with open(conf_file, 'r') as json_file:
        config = json.load(json_file)
        api_key = config['perplexity']

    client = OpenAI(
        api_key=api_key, 
        base_url = "https://api.perplexity.ai" 
    )

    chat_completion = client.chat.completions.create(
        model="sonar-pro", # perplexity
        messages=[
            {
                "role": "user",
                "content": (
                    prompt
                )
            }
        ],
    )
    try:
        citations = getattr(chat_completion, 'citations', [])
        if not citations:
            citations = getattr(chat_completion, 'search_results', [])
    except:
        citations = []
    
    response = chat_completion.choices[0].message.content
    response = _simple_text_formatting(response)
    response = response + "\n\n" + "\n".join(citations)
    save_GPT_response(code, response)

    return response

def _simple_text_formatting(text):
    text = re.sub(r"\*\*(.+?)\*\*", r"[\1]", text)                  # **bold** → [bold]
    text = re.sub(r"^-{3,}\s*$", "\n", text, flags=re.MULTILINE)    # line of 3+ dashes → empty line
    text = re.sub(r"(\n\s*){3,}", "\n\n", text)                     # 3+ empty/whitespace lines → 1 empty line
    return text.strip()

# score_trend, _ = get_score_trend()
# print(score_trend)

# with open('CCA/temp/data.pkl', 'wb') as f:
#     pickle.dump(data_dict, f)

data_dict = post_process(score_trend)

with open('CCA/temp/data.pkl', 'rb') as f:
    data_dict = pickle.load(f)

# print(data_dict)
print(data_dict['select_codelist'])
generate_PPT(data_dict, summary_only = False)
