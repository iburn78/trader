#%% 
# CCA: Company Classification Analysis
from trader.tools.tools import get_df_krx
from trader.data_collection.dc18_Info_DB import get_company_issues, save_GPT_response
import FinanceDataReader as fdr
import pandas as pd
import numpy as np
import os
from html2image import Html2Image
from openai import OpenAI
import json

# ----------------------------------------------------------
# BELOW: Classification Logic
# ----------------------------------------------------------
cv_threshold_prime = 0.4
cv_threshold = 1.0
criteria_dict = {
    'revenue_growth': [(15,7), 0.3], # percent (yoy), # count (e.g., 0.2 = 20% of quarters: to be multiplied by period)
    'revenue_stats': [np.nan, cv_threshold_prime, 0, np.nan], # size
    'opincome_stats': [np.nan, cv_threshold_prime, 0, np.nan], # size
    'opmargin_stats': [(20,10), cv_threshold_prime, 0, np.nan], # percent
    'nopincome_stats': [np.nan, cv_threshold, np.nan, np.nan], # size
    'asset_stats': [np.nan, cv_threshold, 0, np.nan], # size
    'debt_stats': [np.nan, cv_threshold, np.nan, np.nan], # size
    'equity_stats': [np.nan, cv_threshold, 0, np.nan], # size
    'liquid_asset_ratio_stats': [np.nan, cv_threshold, np.nan, np.nan], # percent
    'liquid_debt_ratio_stats': [np.nan, cv_threshold, np.nan, np.nan], # percent
    'debt_to_equity_ratio_stats': [200, cv_threshold, np.nan, np.nan] # percent
}
weight = [5, 2, 2, 5, 1, 1, 1, 1, 1, 1, 1] # weights for each criteria
top_N = 30 # top to add

# ----------------------------------------------------------
# BELOW: Classification defs
# ----------------------------------------------------------

CONF_FILE = os.path.join(os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__)))), 'config/config.json')

def get_qa_db(): 
    pd_ = os.path.dirname(os.path.dirname(os.path.abspath(__file__))) # ..
    qa_db_file = os.path.join(pd_, 'data_collection/data/qa_db.pkl') 
    return pd.read_pickle(qa_db_file)

qa_db = get_qa_db()
df_krx = get_df_krx()

def get_periods(qa_db=qa_db): 
    return [i for i in qa_db.columns if 'Q' in i]

def compare_logic(data_dict, criteria_dict, period:str):
    key = list(data_dict.keys())
    if key != list(criteria_dict.keys()):
        raise ValueError("Key mismatch between data_dict and criteria_dict")

    score = [0] * len(data_dict.keys())
    comment = [''] * len(data_dict.keys())

    # revenue growth
    idx = 0
    data = data_dict[key[idx]]
    crit = criteria_dict[key[idx]] 
    p = int(period.replace('Q',''))
    conditions = [
        data[1] <= int(p*crit[1]), 
    ]
    if all(conditions):
        if data[0] >= crit[0][0]:
            score[idx] = weight[idx]
        elif data[0] >= crit[0][1]:
            score[idx] = round(weight[idx]*(data[0]-crit[0][1])/(crit[0][0]-crit[0][1]), 1)
        else:
            score[idx] = 0
    if data[1] == 0:
        comment[idx] += 'N' # No negative growth'

    # revenue stats, opincome stats, asset stats, equity stats
    for idx in [1, 2, 5, 7]: 
        data = data_dict[key[idx]]
        crit = criteria_dict[key[idx]] 
        conditions = [
            data[1] <= crit[1], # cv
            data[2] >= crit[2], # slope
        ]
        if all(conditions):
            score[idx] = weight[idx]
        if data[3] >= 0:
            comment[idx] += 'A' # Acceration

    # opmargin stats
    for idx in [3]:
        data = data_dict[key[idx]]
        crit = criteria_dict[key[idx]] 
        conditions = [
            data[1] <= crit[1], # cv
            data[2] >= crit[2], # slope
        ]
        if all(conditions):
            if data[0] >= crit[0][0]:
                score[idx] = weight[idx]
            elif data[0] >= crit[0][1]:
                score[idx] = round(weight[idx]*(data[0]-crit[0][1])/(crit[0][0]-crit[0][1]), 1)
            else:
                score[idx] = 0
        if data[3] >= 0:
            comment[idx] += 'A' # Acceration

    # nopincome stats, debt stats, liquid_asset_ratio_stats, liquid_debt_ratio_stats
    for idx in [4, 6, 8, 9]:
        data = data_dict[key[idx]]
        crit = criteria_dict[key[idx]] 
        conditions = [
            data[1] <= crit[1], # cv
        ]
        if all(conditions):
            score[idx] = weight[idx]

    # debt_to_equity_ratio_stats
    for idx in [10]:
        data = data_dict[key[idx]]
        crit = criteria_dict[key[idx]] 
        conditions = [
            data[0] <= crit[0], # mean
            data[1] <= crit[1], # cv
        ]
        if all(conditions):
            score[idx] = weight[idx]
        if data[2] >= 0:
            comment[idx] += 'I' # Increasing
    
    return score, comment

def classify_companies(codelist, criteria_dict, period:str, qa_db=qa_db):
    selected_companies = []
    scores_dict = {}
    for code in codelist:
        data_dict = qa_db.loc[code, period]
        if pd.isna(data_dict) == False:
            scores, comment = compare_logic(data_dict, criteria_dict, period)
            if scores.count(0) == 0:
                selected = True
                selected_companies.append(code)
                # print(code, selected, sum(scores), scores, comment)
            else: 
                selected = False
                pass
            # print(code, selected, sum(scores), scores, comment)
            scores_dict[code] = {
                'selected': selected,
                'score': sum(scores),
                'result': scores,
                'comment': comment
            }
    scores_df = pd.DataFrame.from_dict(scores_dict, orient='index')
    scores_df = scores_df.sort_values(by='score', ascending=False) 
    return selected_companies, scores_df

def get_score_trend(criteria_dict = criteria_dict, df_krx = df_krx, qa_db = qa_db, top_N = top_N):
    codelist = qa_db.index.tolist()
    periods = get_periods(qa_db)

    all_selected_companies = []
    all_scores_dict = {}
    for period in periods:
        selected_companies, scores_df  = classify_companies(codelist, criteria_dict, period)
        all_selected_companies += selected_companies
        all_scores_dict[period] = scores_df

    top_codes = []
    for period in periods: 
        top_codes += all_scores_dict[period].index[:top_N].tolist()

    codelist = list(set(all_selected_companies + top_codes))

    score_trend_df = pd.DataFrame(index = codelist, columns = ['name', 'selected', 'top'+str(top_N)]+periods)
    for code in score_trend_df.index:
        score_trend_df.loc[code, 'name'] = df_krx.loc[code, 'Name']
        selected = ''
        top = ''
        for period in periods:
            if code in all_scores_dict[period].index:
                score_trend_df.loc[code, period] = all_scores_dict[period].loc[code, 'score']
                selected += 'T' if all_scores_dict[period].loc[code, 'selected'] else '-'
                top += 'Y' if code in all_scores_dict[period].index[:top_N] else '-'
            else:
                selected += ' '
                top += ' '
        score_trend_df.loc[code, 'selected'] = selected
        score_trend_df.loc[code, 'top'+str(top_N)] = top

    score_trend_df['avg'] = pd.to_numeric(score_trend_df[periods].mean(axis=1)).round(2)
    score_trend_df = score_trend_df.sort_values(by='avg', ascending=False)
    return score_trend_df, all_scores_dict


# ----------------------------------------------------------
# BELOW: visualize datadict
# ----------------------------------------------------------

def _preprocess_showdata(data_dict):
    data_dict = data_dict.copy()
    rv = data_dict['revenue_growth']
    del data_dict['revenue_growth']

    stats = pd.DataFrame(data_dict, index = ['mean', 'cv', 'slope', 'acc']).T
    stats.index = stats.index.str.replace('_stats', '')
    stats.index = stats.index.str.replace('_ratio', '')
    stats.index = stats.index.str.replace('_', '.')

    # formatting
    stats['slope'] = stats['slope']/stats['mean'].abs()*100
    stats['acc'] = stats['acc']/stats['mean'].abs()*100

    def format_value(val):
        if pd.isna(val):
            return "–"
        return f"{int(val):,}" if abs(val) > 100 else f"{val:.2f}"

    def set_cv_color(val):
        if pd.isna(val):
            return "color: gray"
        if abs(val) >= 0.5:
            return "color: red; font-weight: bold"
        if abs(val) >= 0.2:
            return "color: gray; font-weight: bold"
        else:
            return "color: black; font-weight: bold"
    
    def set_sa_color(val):
        if pd.isna(val):
            return "color: gray"
        if val >= 2:
            return "color: black; font-weight: bold"
        elif val >= 0:
            return "color: gray; font-weight: bold"
        else:
            return "color: red; font-weight: bold" 

    styled = stats[['cv', 'slope', 'acc']].style.format(format_value).map(set_sa_color, subset=['slope', 'acc']).map(set_cv_color, subset=['cv'])
    return rv, stats, styled

def gen_qdata(code, period:str, qa_db = qa_db):
    data_dict = qa_db.loc[code, period]
    if pd.isna(data_dict) or len(data_dict) == 0:
        return "", None
    meta = qa_db.loc[code, 'meta']
    rv, stats, styled = _preprocess_showdata(data_dict)
    p = int(period.replace("Q", ""))

    desc = meta['name'] + ' ' + code + ' / period: ' + period + ' / LQ: ' + meta['last_quarter'] + '\n'
    desc += f'revenue growth:   {rv[0]} %' + '\n'
    desc += f'dip (times):      {rv[1]} ({int(rv[1]/p*100)}%)' + '\n'
    desc += f"op margin:        {stats.loc['opmargin', 'mean']} %" + '\n'
    desc += f"debt to equity:   {stats.loc['debt.to.equity', 'mean']} %" 

    return desc, styled

def gen_summary(code, period:str = None, qa_db = qa_db):
    if period is None:
        periods = get_periods()
    else: 
        periods = [period]
    meta = qa_db.loc[code, 'meta']
    
    summary = []  # collect rows of data across periods
    for period in periods:
        data_dict = qa_db.loc[code, period]
        if pd.isna(data_dict) or len(data_dict) == 0:
            continue
        rv, stats, _ = _preprocess_showdata(data_dict)
        p = int(period.replace("Q", ""))
        row = {
            'period': period,
            'revenue_growth': rv[0],
            'dip': f"{rv[1]} ({int(rv[1]/p*100)}%)",
            'opmargin': stats.loc['opmargin', 'mean'],
            'debt.to.equity': stats.loc['debt.to.equity', 'mean']
        }
        summary.append(row)

    if len(summary) == 0:
        # print('No data found for any period.')
        return

    summary_df = pd.DataFrame(summary).set_index('period')

    def format_cell(val):
        if pd.isna(val):
            return "–"
        return f"{val:.2f}" if isinstance(val, float) else str(val)

    def highlight_negative(val):
        if isinstance(val, float) and val < 0:
            return "color: red"
        return "color: black"

    styled = summary_df.style.format(format_cell).map(highlight_negative)
    desc = f"{meta['name']} ({code})" + '\n' + f"Latest Quarter: {meta['last_quarter']}"

    return desc, styled
    
def gen_data_in_html(code, qa_db = qa_db): 
    summary_desc, summary_df = gen_summary(code)
    summary_desc = summary_desc.replace('\n', '<br>')
    summary_df = summary_df.to_html()

    qdata_desc_list = []
    qdata_df_list = []
    for period in get_periods():
        qdata_desc, qdata_df = gen_qdata(code, period)
        if qdata_df != None: 
            qdata_desc_list.append(qdata_desc.replace('\n', '<br>'))
            qdata_df_list.append(qdata_df.to_html())
    
    left_html = summary_desc + "<br>" + summary_df + "<br>"
    center_html = ''
    right_html = ''
    for a, b in zip(qdata_desc_list[:1], qdata_df_list[:1]):
        left_html += a + "<br>" + b + "<br>"
    for a, b in zip(qdata_desc_list[1:3], qdata_df_list[1:3]):
        center_html += a + "<br>" + b + "<br>"
    for a, b in zip(qdata_desc_list[3:], qdata_df_list[3:]):
        right_html += a + "<br>" + b + "<br>"

    full_html = f"""
    <html>
    <head>
    <style>
        body {{ font-family: Arial; font-size: 14px; margin: 0; padding: 1px; background: white; }}
        table {{ border-collapse: collapse; }}
        th, td {{ border: 1px solid #ccc; padding: 2px 4px; text-align: center;}}
    </style>
    </head>
    <body>
    <div style="display: flex;">
        <div style="flex: 1; padding: 10px; border: 1px solid #ccc;">
            <p>{left_html}</p>
        </div>
        <div style="flex: 1; padding: 10px; border: 1px solid #ccc;">
            <p>{center_html}</p>
        </div>
        <div style="flex: 1; padding: 10px; border: 1px solid #ccc;">
            <p>{right_html}</p>
        </div>
    </div>
    </body>
    </html>
    """
    cd_ = os.path.dirname(os.path.abspath(__file__))
    temp_html = os.path.join(cd_, 'CCA/temp/temp.html')
    img_path = os.path.join(cd_, 'CCA/temp/')
    filename = 'temp.png'
    with open("CCA/temp/temp.html", "w", encoding="utf-8") as f:
        f.write(full_html)  

    hti = Html2Image(output_path=img_path)
    hti.screenshot(html_file=temp_html, save_as=filename, size=(1400, 900))
    os.remove(temp_html)

    return os.path.join(img_path, filename)

# ----------------------------------------------------------
# BELOW: PPT Generation
# ----------------------------------------------------------

import io
from trader.tools.tools import get_main_financial_reports_db, get_quarterly_data, get_price_db, get_outshare_db, prev_quarter_str
from trader.tools.tools import plot_company_financial_summary2
from trader.data_collection.dc17_QuarterlyAnalysisDB import MAX_QUARTERS
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches

fr_db = get_main_financial_reports_db()
pr_db = get_price_db()
outshare_DB = get_outshare_db()
target_account = 'net_income'

# this is a different version of L4_addition function from anaysis project
def L4_rolling_addition(fh, target_account=target_account):
    # Add L4 key account to the financial report database.
    new_row = {'account':'L4_'+target_account }
    target_row = fh.loc[target_account]

    for i in range(3, len(fh.columns)):
        previous_4_quarters = fh.columns[i-3:i+1]
        rolling_sum = target_row[previous_4_quarters].sum()
        new_row[fh.columns[i]] = rolling_sum

    new_row_df = pd.DataFrame([new_row]).set_index('account')
    fh = pd.concat([fh, new_row_df])
    return fh 

# get price, marcap, PER, PBR
def get_market_performance_db(code, fr_db=fr_db, pr_db=pr_db, outshare_DB=outshare_DB, target_account=target_account, MAX_QUARTERS=MAX_QUARTERS):
    fh = get_quarterly_data(code, fr_db, native=True)
    fh = L4_rolling_addition(fh, target_account)
    init_loc = fh.columns.get_loc(fh.loc['L4_'+ target_account].first_valid_index())
    start_loc = max(max(len(fh.columns) - MAX_QUARTERS, 0), init_loc)
    start_day  = pd.Period(fh.columns[start_loc][:5].replace('_', ''), freq='Q').start_time.strftime('%Y-%m-%d')

    # market performance
    mp_db = pd.DataFrame()
    mp_db['price'] = pr_db[code].loc[pr_db.index >= start_day]
    mp_db.index = pd.to_datetime(mp_db.index)
    shares = outshare_DB[code].iloc[-1]
    mp_db['marcap'] = mp_db['price'] * shares 
    r_ = fh.loc['L4_' + target_account].dropna()
    latest_l4 = r_.iloc[-1] if not r_.empty else np.nan
    mp_db['L4'] = mp_db.index.map(lambda d: fh.at['L4_'+target_account, prev_quarter_str(d)] if prev_quarter_str(d) in fh.columns else latest_l4)
    mp_db['PER'] = mp_db['marcap'] / mp_db['L4']

    b_ = fh.loc['equity'].dropna()
    latest_b = b_.iloc[-1] if not b_.empty else np.nan
    mp_db['equity'] = mp_db.index.map(lambda d: fh.at['equity', prev_quarter_str(d)] if prev_quarter_str(d) in fh.columns else latest_b)
    mp_db['PBR'] = mp_db['marcap'] / mp_db['equity']
    mp_db = mp_db.apply(pd.to_numeric, errors='coerce')

    return mp_db


def mp_plot(mp_db, columns=['price', 'PER', 'PBR']):
    fig, axes = plt.subplots(len(columns), 1, figsize=(12, 4 * len(columns)), sharex=True)
    if len(columns) == 1:
        axes = [axes]
    for ax, col in zip(axes, columns):
        mp_db[col].plot(ax=ax, title=col, grid=True, fontsize=12)
    plt.tight_layout()
    img_stream = io.BytesIO()
    plt.savefig(img_stream, format='png', bbox_inches='tight')
    plt.close(fig)
    img_stream.seek(0)
    return img_stream
    # plt.show()
    # plt.savefig("----.png")


def generate_PPT(score_trend, codelist=None, fr_db=fr_db, pr_db=pr_db, outshare_DB=outshare_DB, df_krx=df_krx, topN = 100):
    cd_ = os.path.dirname(os.path.abspath(__file__)) # .
    today_str = pd.Timestamp.today().strftime('%Y-%m-%d_%H%M')
    CCA_template = os.path.join(cd_, 'CCA/CCA_template.pptx')
    CCA_result = os.path.join(cd_, f'CCA/CCA_result_{today_str}.pptx')
    prs = Presentation(CCA_template)
    periods = get_periods(qa_db)

    for code in score_trend.index[:topN]:
        if codelist != None and code not in codelist: 
            continue
        print("processing", code)
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        mp_db = get_market_performance_db(code, fr_db, pr_db, outshare_DB)
        img_stream = mp_plot(mp_db)
        slide.shapes.add_picture(img_stream, Inches(0.1), Inches(0.5), height=Inches(6))
        for ph in slide.placeholders: 
            if ph.name == 'Text Placeholder 1':
                ph.text = df_krx.loc[code, 'Name'] + f" ({code})"
            if ph.name == 'Text Placeholder 2':
                ph.text = today_str
            if ph.name == 'Text Placeholder 3':
                txt = score_trend.loc[[code]][periods+['avg']].to_string(index=False)
                rank = str(score_trend.index.get_loc(code) + 1)
                ph.text = 'rank:' + rank + ', select:' + score_trend.loc[code, 'selected'] + ', top30:' + score_trend.loc[code, 'top30'] + ', MCap:' + str(df_krx.at[code, 'Marcap'] // 10**8) + ', PER:' + str(round(mp_db['PER'].iloc[-1], 2)) + '\n' + txt

        slide = prs.slides.add_slide(prs.slide_layouts[1])
        ph = next(iter(slide.placeholders)) 
        ph.text = openai_command(code)

        slide = prs.slides.add_slide(prs.slide_layouts[2])
        img_stream = plot_company_financial_summary2(fr_db, pr_db, code, None) 
        slide.shapes.add_picture(img_stream, Inches(0.1), 0, height=prs.slide_height)

        slide = prs.slides.add_slide(prs.slide_layouts[2])
        img_path = gen_data_in_html(code)
        slide.shapes.add_picture(img_path, 0, Inches(0.1), width=prs.slide_width)
        os.remove(img_path)

    prs.save(CCA_result)
    print("PPT generation completed...")

# open ai original
# def openai_command(company_name, conf_file=CONF_FILE):  
#     with open(conf_file, 'r') as json_file:
#         config = json.load(json_file)
#         api_key = config['openai']

#     client = OpenAI(
#         api_key=api_key
#     )
#     content_command = f"{company_name} is a listed company in Korea. Give me brief description (1) what is this company's business, (2) why this company's performance is good for the last quarters, (3) what are the key 3 issues ths company is facing?"
#     format_command = "Answer in Korean language, and use only plain text. Total length of answers should not exceed 500 words."

#     chat_completion = client.chat.completions.create(
#         model="gpt-4o-mini", 
#         messages=[
#             {
#                 "role": "user",
#                 "content": (
#                     content_command + " " + format_command
#                 )
#             }
#         ]
#     )
#     return chat_completion.choices[0].message.content

# perplexity
def openai_command(code, conf_file=CONF_FILE):
    company_name = df_krx.loc[code, "Name"]
    with open(conf_file, 'r') as json_file:
        config = json.load(json_file)
        api_key = config['perplexity']

    client = OpenAI(
        api_key=api_key, 
        base_url = "https://api.perplexity.ai" 
    )

    content_command = (
        f"{company_name}({code}) is a listed company in Korea. Provide a detailed and specific explanation of "
        "(1) what this company's core business is, "
        "(2) why the company has shown strong performance in recent quarters, and "
        "(3) what the top 3 key issues the company is currently facing. "
    )

    format_command = (
        "Answer in Korean using plain text. "
        "Avoid vague or high-level generalizations. "
        "The total length must not exceed 700 words. "
    )

    info, prev, date_prev = get_company_issues(code)
    if info:
        info_command = (
            "The following is additional information. "
            "Cross-check its validity and use it only if it is important and accurate.\n"
        )
        info = info_command + info

    if prev:
        prev_command = (
            f"The following is your previous response to the identical query about the company on {date_prev}. "
            "Refer to it, but update your answer with any new developments since then.\n"
        )
        prev = prev_command + prev

    prompt = f"{content_command}\n{format_command}\n\n{info}\n\n{prev}"

    chat_completion = client.chat.completions.create(
        model="sonar-pro", # perplexity
        messages=[
            {
                "role": "user",
                "content": (
                    prompt
                )
            }
        ] 
    )
    try:
        citations = getattr(chat_completion, 'citations', [])
        if citations is []:
            citations = getattr(chat_completion, 'search_results', [])
    except:
        citations = []
    
    response = chat_completion.choices[0].message.content + "\n\n" + "\n".join(citations)
    save_GPT_response(code, response)

    return response
