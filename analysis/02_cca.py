#%% 
# ----------------------------------------------------------
# CCA: Company Classification Analysis
# CCA Post processing and PPT Generation
# ----------------------------------------------------------
# PART 1: cca_dict creation
# - all operation is local
# ----------------------------------------------------------
import pickle
from datetime import datetime, timedelta
import io, os
import pandas as pd
from tqdm import tqdm
from trader.tools.dc_tools import get_main_financial_reports_db, get_quarterly_data, plot_company_financial_summary
from trader.tools.cca_tools import get_score_trend, get_periods, L4_rolling_addition, prev_quarter_str
from trader.tools.cca_tools import df_krx, qa_db, top_N

fr_db = get_main_financial_reports_db()
pd_ = os.path.dirname(os.path.dirname(os.path.abspath(__file__))) # ..
pr_db_file = os.path.join(pd_, 'data_collect/data/price_DB.feather') 
pr_db = pd.read_feather(pr_db_file)
cca_dict_file = os.path.join(pd_, 'data_collect/cca/temp/cca_dict.pkl')

target_account = 'net_income' # for some companies, api does not provide net income data
target_account = 'operating_income' # PER is then based on operating income, could be different from other PER publications

# mp_db: just to get price, marcap, PER, PBR
def get_market_performance_db(code, fr_db=fr_db, pr_db=pr_db, target_account=target_account, MAX_QUARTERS=40): # e.g., 10 years or use MAX_QUARTERS in qa_db generation
    fh = get_quarterly_data(code, fr_db, native=True)
    fh = L4_rolling_addition(fh, target_account)
    try:
        init_loc = fh.columns.get_loc(fh.loc['L4_'+ target_account].first_valid_index())
        start_loc = max(max(len(fh.columns) - MAX_QUARTERS, 0), init_loc)
    except:
        start_loc = 0
    start_day  = pd.Period(fh.columns[start_loc][:5].replace('_', ''), freq='Q').start_time.strftime('%Y-%m-%d')

    # market performance
    mp_db = pd.DataFrame()
    mp_db['price'] = pr_db[code].loc[pr_db.index >= start_day]
    mp_db.index = pd.to_datetime(mp_db.index)
    shares = df_krx.at[code, 'Stocks']
    mp_db['marcap'] = mp_db['price'] * shares 
    r_ = fh.loc['L4_' + target_account].dropna()
    latest_l4 = r_.iloc[-1] if not r_.empty else np.nan
    mp_db['L4'] = mp_db.index.map(lambda d: fh.at['L4_'+target_account, prev_quarter_str(d)] if prev_quarter_str(d) in fh.columns else latest_l4)
    mp_db['PER'] = mp_db['marcap'] / mp_db['L4']

    b_ = fh.loc['equity'].dropna()
    latest_b = b_.iloc[-1] if not b_.empty else np.nan
    mp_db['equity'] = mp_db.index.map(lambda d: fh.at['equity', prev_quarter_str(d)] if prev_quarter_str(d) in fh.columns else latest_b)
    mp_db['PBR'] = mp_db['marcap'] / mp_db['equity']
    mp_db = mp_db.apply(pd.to_numeric, errors='coerce')

    return mp_db

def _get_late_prices(codelist, days=7):
    start_date = (datetime.today() - timedelta(days=days)).strftime('%Y-%m-%d')
    res = pd.DataFrame()

    for code in codelist:
        df = pr_db.loc[pr_db.index >= start_date, code] # df.name == code (remains as the column name)
        res = pd.concat([res, df], axis=1)

    res = res.T
    res.columns = pd.to_datetime(res.columns).strftime("%m-%d")
    return res  # transpose to have stock codes as index and dates as columns

def _get_codelist_summary(codelist, days=7, df_krx=df_krx):
    last_prices = _get_late_prices(codelist, days)
    pr_changes = (last_prices.pct_change(axis=1) * 100).round(2).dropna(axis=1)
    pr_changes.insert(0, 'name', '')
    for code in pr_changes.index:
        pr_changes.loc[code, 'name'] = df_krx.loc[code, "Name"]
        pr_changes.loc[code, 'price'] = last_prices.loc[code, last_prices.columns[-1]]
        pr_changes.loc[code, 'marcap'] = df_krx.loc[code, "Marcap"]/10**8
    pr_changes['marcap'] = pr_changes['marcap'].round().astype(int)
    return pr_changes

def _post_process(score_trend, qa_db=qa_db, fr_db=fr_db, pr_db=pr_db, df_krx=df_krx, top_N = top_N):
    cca_dict = {}
    mp_db_dict = {}
    periods = get_periods(qa_db)

    cls = _get_codelist_summary(score_trend.index)
    pbar = tqdm(score_trend.index, desc="score_trend items")
    for code in pbar:
        mp_db = get_market_performance_db(code, fr_db, pr_db)
        mp_db_dict[code] = mp_db
        cls.loc[code, 'PER'] = round(mp_db['PER'].iloc[-1], 2)
        rank = str(score_trend.index.get_loc(code) + 1)
        cls.loc[code, 'rank'] = rank 
        txt = score_trend.loc[[code]][periods+['avg']].to_string(index=False)
        cls.loc[code, 'note'] = 'rank:' + rank + ', select:' + score_trend.loc[code, 'selected'] + f', top{str(top_N)}:' + score_trend.loc[code, f'top{str(top_N)}'] + ', MCap:' + str(df_krx.at[code, 'Marcap'] // 10**8) + ', PER:' + str(round(mp_db['PER'].iloc[-1], 2)) + '\n' + txt

    cca_dict['created'] = datetime.today().strftime('%Y-%m-%d')
    cca_dict['codelist_summary'] = cls
    cca_dict['mp_db_dict'] = mp_db_dict

    # -------------------------------------------------------
    # additional codelist selection logic  
    # -------------------------------------------------------
    PER_limit = 1000
    scodelist = cls.loc[(cls['PER'] > 0) & (cls['PER'] < PER_limit)].index
    # -------------------------------------------------------
    cca_dict['select_codelist'] = scodelist
    cca_dict['select_codelist_summary'] = cls.loc[cls.index.isin(scodelist)].drop('note', axis=1)

    return cca_dict

def get_cca_dict(cca_dict_file, force_recreate=False):
    os.makedirs(os.path.dirname(cca_dict_file), exist_ok=True)

    _today = datetime.today().strftime('%Y-%m-%d')
    needs_refresh = force_recreate or not os.path.exists(cca_dict_file)

    if not needs_refresh:
        try:
            with open(cca_dict_file, 'rb') as f:
                cca_dict = pickle.load(f)

            needs_refresh = cca_dict.get('created') != _today

        except Exception:
            # Corrupt or unreadable cache; rebuild it.
            needs_refresh = True

    if needs_refresh:
        score_trend, _ = get_score_trend()
        cca_dict = _post_process(score_trend)

        with open(cca_dict_file, 'wb') as f:
            pickle.dump(cca_dict, f)

    return cca_dict

# ----------------------------------------------------------
# PART 2: PPT generation
# - uses LLM: can be non-local
# - input: cca_dict, fr_db, pr_db, top_N
# ----------------------------------------------------------
import re
from openai import OpenAI
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches
from trader.tools.cca_tools import styled_df_to_image, gen_data_in_html, get_prev_response_and_issues, save_LLM_response

temp_path = os.path.join(pd_, 'data_collect/cca/temp/')
os.makedirs(temp_path, exist_ok=True)

CCA_template = os.path.join(pd_, 'data_collect/cca/util/CCA_template.pptx')

today_hm = pd.Timestamp.today().strftime('%Y-%m-%d_%H%M')
CCA_result = os.path.join(pd_, f'data_collect/cca/CCA_result_{today_hm}.pptx')

def _mp_plot(mp_db, columns=['price', 'PER', 'PBR']):
    fig, axes = plt.subplots(len(columns), 1, figsize=(12, 4 * len(columns)), sharex=True)
    if len(columns) == 1:
        axes = [axes]
    for ax, col in zip(axes, columns):
        mp_db[col].plot(ax=ax, title=col, grid=True, fontsize=12)
    plt.tight_layout()
    img_stream = io.BytesIO()
    plt.savefig(img_stream, format='png', bbox_inches='tight')
    img_stream.seek(0)
    # plt.show()
    plt.close(fig)
    return img_stream

def generate_PPT(cca_dict, fr_db=fr_db, pr_db=pr_db, summary_only=False, top_N: int = top_N):
    prs = Presentation(CCA_template)

    # Codelist summary page 
    slide = prs.slides.add_slide(prs.slide_layouts[2])
    img_path = styled_df_to_image(cca_dict['select_codelist_summary'], temp_path)
    slide.shapes.add_picture(img_path, 0, Inches(0.1), width=prs.slide_width*0.55)
    os.remove(img_path)

    # Per code pages 
    if not summary_only: 
        if top_N is None: top_N = len(cca_dict['select_codelist'])
        pbar = tqdm(cca_dict['select_codelist'][:top_N])
        for code in pbar:
            pbar.set_postfix(code=code)
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            img_stream = _mp_plot(cca_dict['mp_db_dict'][code])
            slide.shapes.add_picture(img_stream, Inches(0.1), Inches(0.5), height=Inches(6))
            for ph in slide.placeholders: 
                if ph.name == 'Text Placeholder 1':
                    ph.text = cca_dict['codelist_summary'].loc[code, 'name']
                if ph.name == 'Text Placeholder 2':
                    ph.text = today_hm
                if ph.name == 'Text Placeholder 3':
                    ph.text = cca_dict['codelist_summary'].loc[code, 'note']

            slide = prs.slides.add_slide(prs.slide_layouts[1])
            ph = next(iter(slide.placeholders)) 
            ph.text = LLM_request(code)

            slide = prs.slides.add_slide(prs.slide_layouts[2])
            img_stream = plot_company_financial_summary(fr_db, pr_db, code, None) 
            slide.shapes.add_picture(img_stream, Inches(0.1), 0, height=prs.slide_height)

            slide = prs.slides.add_slide(prs.slide_layouts[2])
            img_path = gen_data_in_html(code, temp_path)
            slide.shapes.add_picture(img_path, 0, Inches(0.1), width=prs.slide_width)
            os.remove(img_path)

    prs.save(CCA_result)
    print("PPT generation completed...")

def _simple_text_formatting(text):
    text = re.sub(r"\*\*(.+?)\*\*", r"[\1]", text)                  # **bold** → [bold]
    text = re.sub(r"^-{3,}\s*$", "\n", text, flags=re.MULTILINE)    # line of 3+ dashes → empty line
    text = re.sub(r"(\n\s*){3,}", "\n\n", text)                     # 3+ empty/whitespace lines → 1 empty line
    return text.strip()

def LLM_request(code):
    company_name = df_krx.loc[code, "Name"]

    content_command = (
        f"{company_name}({code}) is a listed company in Korea. Provide a detailed and specific explanation of "
        "(1) what this company's core business is, "
        "(2) why the company has shown strong performance in recent quarters, and "
        "(3) what the top 3 key issues the company is currently facing. "
    )

    format_command = (
        "Answer in Korean using plain text. "
        "Avoid vague or high-level generalizations. "
        "The total length must not exceed 700 words. "
    )

    issues, prev_response, date_prev = get_prev_response_and_issues(code)

    # if the prev gpt response is saved within 7 days, pass
    if datetime.today() - datetime.strptime(date_prev, '%Y-%m-%d') < timedelta(days=7):
        if prev_response:
            return prev_response

    if issues:
        issues_command = (
            "The following is additional information. "
            "Cross-check its validity and use it only if it is important and accurate.\n"
        )
        issues = issues_command + issues

    if prev_response:
        prev_command = (
            f"The following is your previous response to the identical query about the company on {date_prev}. "
            "Refer to it, but update your answer with any new developments since then.\n"
        )
        prev_info = prev_command + prev_response
    
    extras = "\n\n".join(filter(None, [issues, prev_info]))
    prompt = f"{content_command}\n{format_command}\n\n{extras}"

    client = OpenAI(
        api_key= "dummy", 
        base_url="http://localhost:11434/v1", # ollama
        # base_url = "https://api.perplexity.ai" # perplexity API may contain reference info, separately
    )

    # ----------------------------------
    # decommand the following to use LLM
    # ----------------------------------
    # chat_completion = client.chat.completions.create(
    #     model="gemma4", 
    #     # model="sonar-pro", # perplexity
    #     messages=[
    #         {
    #             "role": "user",
    #             "content": (
    #                 prompt
    #             )
    #         }
    #     ],
    # )
    
    # response = chat_completion.choices[0].message.content
    # response = _simple_text_formatting(response)
    # save_LLM_response(code, response)

    # ----------------------------------
    # remove the following to enable LLM
    # ----------------------------------
    response = prev_response

    return response

if __name__ == "__main__":
    cca_dict = get_cca_dict(cca_dict_file, force_recreate=False)
    generate_PPT(cca_dict, summary_only=False)
