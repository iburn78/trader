#%% 
# ----------------------------------------------------------
# CCA: Company Classification Analysis
# CCA Post processing and PPT Generation
# ----------------------------------------------------------
import pickle
from openai import OpenAI
import re
from datetime import datetime, timedelta
import FinanceDataReader as fdr
import io, os
import pandas as pd
from trader.tools.dc_tools import get_main_financial_reports_db, plot_company_financial_summary
from trader.tools.cca_tools import get_score_trend, get_quarterly_data, prev_quarter_str, get_periods, styled_df_to_image, get_company_issues, save_GPT_response
from trader.tools.cca_tools import df_krx, qa_db, top_N
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches

fr_db = get_main_financial_reports_db()
pd_ = os.path.dirname(os.path.dirname(os.path.abspath(__file__))) # ..
pr_db_file = os.path.join(pd_, 'data_collect/data/price_DB.feather') 
pr_db = pd.read_feather(pr_db_file)

target_account = 'net_income'

def L4_rolling_addition(fh, target_account=target_account):
    # Add L4 target account to the financial report database.
    new_row = {'account':'L4_'+target_account }
    target_row = fh.loc[target_account]

    for i in range(3, len(fh.columns)):
        previous_4_quarters = fh.columns[i-3:i+1]
        rolling_sum = target_row[previous_4_quarters].sum()
        new_row[fh.columns[i]] = rolling_sum

    new_row_df = pd.DataFrame([new_row]).set_index('account')
    fh = pd.concat([fh, new_row_df])
    return fh 

# get price, marcap, PER, PBR
def get_market_performance_db(code, fr_db=fr_db, pr_db=pr_db, target_account=target_account, MAX_QUARTERS=40): # e.g., 10 years or use MAX_QUARTERS in qa_db generation
    fh = get_quarterly_data(code, fr_db, native=True)
    fh = L4_rolling_addition(fh, target_account)
    try:
        init_loc = fh.columns.get_loc(fh.loc['L4_'+ target_account].first_valid_index())
        start_loc = max(max(len(fh.columns) - MAX_QUARTERS, 0), init_loc)
    except:
        start_loc = 0
    start_day  = pd.Period(fh.columns[start_loc][:5].replace('_', ''), freq='Q').start_time.strftime('%Y-%m-%d')

    # market performance
    mp_db = pd.DataFrame()
    mp_db['price'] = pr_db[code].loc[pr_db.index >= start_day]
    mp_db.index = pd.to_datetime(mp_db.index)
    shares = df_krx.at[code, 'Stocks']
    mp_db['marcap'] = mp_db['price'] * shares 
    r_ = fh.loc['L4_' + target_account].dropna()
    latest_l4 = r_.iloc[-1] if not r_.empty else np.nan
    mp_db['L4'] = mp_db.index.map(lambda d: fh.at['L4_'+target_account, prev_quarter_str(d)] if prev_quarter_str(d) in fh.columns else latest_l4)
    mp_db['PER'] = mp_db['marcap'] / mp_db['L4']

    b_ = fh.loc['equity'].dropna()
    latest_b = b_.iloc[-1] if not b_.empty else np.nan
    mp_db['equity'] = mp_db.index.map(lambda d: fh.at['equity', prev_quarter_str(d)] if prev_quarter_str(d) in fh.columns else latest_b)
    mp_db['PBR'] = mp_db['marcap'] / mp_db['equity']
    mp_db = mp_db.apply(pd.to_numeric, errors='coerce')

    return mp_db

def mp_plot(mp_db, columns=['price', 'PER', 'PBR']):
    fig, axes = plt.subplots(len(columns), 1, figsize=(12, 4 * len(columns)), sharex=True)
    if len(columns) == 1:
        axes = [axes]
    for ax, col in zip(axes, columns):
        mp_db[col].plot(ax=ax, title=col, grid=True, fontsize=12)
    plt.tight_layout()
    img_stream = io.BytesIO()
    plt.savefig(img_stream, format='png', bbox_inches='tight')
    plt.close(fig)
    img_stream.seek(0)
    return img_stream
    # plt.show()
    # plt.savefig("----.png")

def get_late_prices(codelist, days=7):
    start_date = (datetime.today() - timedelta(days=days)).strftime('%Y-%m-%d')
    res = pd.DataFrame()

    for code in codelist:
        df = fdr.DataReader(code, start_date)['Close']
        df.name = code  # name the Series as the stock code
        res = pd.concat([res, df], axis=1)

    res = res.T
    res.columns = pd.to_datetime(res.columns).strftime("%m-%d")
    return res  # transpose to have stock codes as index and dates as columns

def get_codelist_summary(codelist, days=7, df_krx=df_krx):
    last_prices = get_late_prices(codelist, days)
    pr_changes = (last_prices.pct_change(axis=1) * 100).round(2).dropna(axis=1)
    pr_changes.insert(0, 'name', '')
    for code in pr_changes.index:
        pr_changes.loc[code, 'name'] = df_krx.loc[code, "Name"]
        pr_changes.loc[code, 'price'] = last_prices.loc[code, last_prices.columns[-1]]
        pr_changes.loc[code, 'marcap'] = df_krx.loc[code, "Marcap"]/10**8
    pr_changes['marcap'] = pr_changes['marcap'].round().astype(int)
    return pr_changes

def post_process(score_trend, qa_db=qa_db, fr_db=fr_db, pr_db=pr_db, df_krx=df_krx, top_N = top_N):
    data_dict = {}
    mp_db_dict = {}
    periods = get_periods(qa_db)

    cls = get_codelist_summary(score_trend.index)
    for code in score_trend.index: 
        print("processing", code)
        mp_db = get_market_performance_db(code, fr_db, pr_db)
        mp_db_dict[code] = mp_db
        cls.loc[code, 'PER'] = round(mp_db['PER'].iloc[-1], 2)
        rank = str(score_trend.index.get_loc(code) + 1)
        cls.loc[code, 'rank'] = rank 
        txt = score_trend.loc[[code]][periods+['avg']].to_string(index=False)
        cls.loc[code, 'note'] = 'rank:' + rank + ', select:' + score_trend.loc[code, 'selected'] + f', top{str(top_N)}:' + score_trend.loc[code, f'top{str(top_N)}'] + ', MCap:' + str(df_krx.at[code, 'Marcap'] // 10**8) + ', PER:' + str(round(mp_db['PER'].iloc[-1], 2)) + '\n' + txt

    data_dict['codelist_summary'] = cls
    data_dict['mp_db_dict'] = mp_db_dict

    # codelist selection logic ------------------
    PER_limit = 1000
    scodelist = cls.loc[(cls['PER'] > 0) & (cls['PER'] < PER_limit)].index
    data_dict['select_codelist'] = scodelist
    data_dict['select_codelist_summary'] = cls.loc[cls.index.isin(scodelist)].drop('note', axis=1)

    return data_dict

def generate_PPT(data_dict, fr_db=fr_db, pr_db=pr_db, summary_only = False, top_N: int = top_N):

    cd_ = os.path.dirname(os.path.abspath(__file__)) # .
    CCA_folder = 'CCA'
    today_str = pd.Timestamp.today().strftime('%Y-%m-%d_%H%M')
    CCA_template = os.path.join(cd_, CCA_folder, 'CCA_template.pptx')
    CCA_result = os.path.join(cd_, CCA_folder, f'CCA_result_{today_str}.pptx')
    prs = Presentation(CCA_template)

    # Codelist summary page 
    slide = prs.slides.add_slide(prs.slide_layouts[2])
    img_path = styled_df_to_image(data_dict['select_codelist_summary'])
    slide.shapes.add_picture(img_path, 0, Inches(0.1), width=prs.slide_width*0.55)
    os.remove(img_path)

    # Per code pages 
    if not summary_only: 
        if top_N is None: top_N = len(data_dict['select_codelist'])
        for code in data_dict['select_codelist'][:top_N]:
            print("slide generating for", code)
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            img_stream = mp_plot(data_dict['mp_db_dict'][code])
            slide.shapes.add_picture(img_stream, Inches(0.1), Inches(0.5), height=Inches(6))
            for ph in slide.placeholders: 
                if ph.name == 'Text Placeholder 1':
                    ph.text = data_dict['codelist_summary'].loc[code, 'name']
                if ph.name == 'Text Placeholder 2':
                    ph.text = today_str
                if ph.name == 'Text Placeholder 3':
                    ph.text = data_dict['codelist_summary'].loc[code, 'note']

            slide = prs.slides.add_slide(prs.slide_layouts[1])
            ph = next(iter(slide.placeholders)) 
            ph.text = LLM_request(code)

            slide = prs.slides.add_slide(prs.slide_layouts[2])
            img_stream = plot_company_financial_summary(fr_db, pr_db, code, None) 
            slide.shapes.add_picture(img_stream, Inches(0.1), 0, height=prs.slide_height)

            slide = prs.slides.add_slide(prs.slide_layouts[2])
            img_path = gen_data_in_html(code)
            slide.shapes.add_picture(img_path, 0, Inches(0.1), width=prs.slide_width)
            os.remove(img_path)

    prs.save(CCA_result)
    print("PPT generation completed...")

def LLM_request(code):
    company_name = df_krx.loc[code, "Name"]

    content_command = (
        f"{company_name}({code}) is a listed company in Korea. Provide a detailed and specific explanation of "
        "(1) what this company's core business is, "
        "(2) why the company has shown strong performance in recent quarters, and "
        "(3) what the top 3 key issues the company is currently facing. "
    )

    format_command = (
        "Answer in Korean using plain text. "
        "Avoid vague or high-level generalizations. "
        "The total length must not exceed 700 words. "
    )

    info, prev, date_prev = get_company_issues(code)

    # if the prev gpt response is saved within 7 days, pass
    if datetime.today() - datetime.strptime(date_prev, '%Y-%m-%d') < timedelta(days=7):
        if prev:
            return prev

    if info:
        info_command = (
            "The following is additional information. "
            "Cross-check its validity and use it only if it is important and accurate.\n"
        )
        info = info_command + info

    if prev:
        prev_command = (
            f"The following is your previous response to the identical query about the company on {date_prev}. "
            "Refer to it, but update your answer with any new developments since then.\n"
        )
        prev = prev_command + prev

    extras = "\n\n".join(filter(None, [info, prev]))
    prompt = f"{content_command}\n{format_command}\n\n{extras}"

    client = OpenAI(
        api_key= "dummy", 
        # base_url = "https://api.perplexity.ai" 
        base_url="http://localhost:11434/v1", # ollama
    )

    chat_completion = client.chat.completions.create(
        # model="sonar-pro", # perplexity
        model="gemma4", 
        messages=[
            {
                "role": "user",
                "content": (
                    prompt
                )
            }
        ],
    )
    try:
        citations = getattr(chat_completion, 'citations', [])
        if not citations:
            citations = getattr(chat_completion, 'search_results', [])
    except:
        citations = []
    
    response = chat_completion.choices[0].message.content
    response = _simple_text_formatting(response)
    response = response + "\n\n" + "\n".join(citations)
    save_GPT_response(code, response)

    return response

def _simple_text_formatting(text):
    text = re.sub(r"\*\*(.+?)\*\*", r"[\1]", text)                  # **bold** → [bold]
    text = re.sub(r"^-{3,}\s*$", "\n", text, flags=re.MULTILINE)    # line of 3+ dashes → empty line
    text = re.sub(r"(\n\s*){3,}", "\n\n", text)                     # 3+ empty/whitespace lines → 1 empty line
    return text.strip()

data_dict_file = os.path.join(pd_, 'data_collect/cca/temp/data_dict.pkl')
def get_data_dict(re_create=False, data_dict_file=data_dict_file):
    
    if re_create or not os.path.exists(data_dict_file):
        score_trend, _ = get_score_trend()
        data_dict = post_process(score_trend)
        with open(data_dict_file, 'wb') as f:
            pickle.dump(data_dict, f)
    else:
        with open(data_dict_file, 'rb') as f:
            data_dict = pickle.load(f)
    return data_dict

if __name__ == "__main__":
    data_dict = get_data_dict(re_create=False)
    generate_PPT(data_dict, summary_only = False)

